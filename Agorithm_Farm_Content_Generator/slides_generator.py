import json
import re
import os
import io

# Import LangChain components
from langchain_openai import ChatOpenAI                     

# 📊 Generate PowerPoint 
from pptx import Presentation
import traceback # For detailed error logging
from pptx.enum.chart import XL_LABEL_POSITION, XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches 
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor # Optional: for text formatting
from tqdm import tqdm # For progress bar
from PIL import Image
from image_search import search_and_generate_images, load_image_captions
import streamlit as st
from st_files_connection import FilesConnection


# Display the PowerPoint presentation
import tempfile
import subprocess
import traceback # For detailed error logging


# Load text data from a preprocessed JSONL file
# Define the path to the processed text file (two levels up from current directory)
script_dir = os.path.dirname(os.path.abspath(__file__))
base_dir = os.path.abspath(os.path.join(script_dir, ".."))
document_path = os.path.join(base_dir, "Data", "Processed", "Text")
ppt_template_path = os.path.join(base_dir, "Data", "PPT_01.pptx")

# Model & Generation Config
# LLM_MODEL = "gpt-4o-mini"
IMAGE_MODEL = "dall-e-3"
IMAGE_SIZE = "1024x1024"
IMAGE_QUALITY = "standard"


# 🔄 Generate structured slide content using LLM
def generate_slide_content(final_prompt, num_slides=10, api_key=None, model_name="gpt-4o-mini-2024-07-18"):
    llm = ChatOpenAI(
        model=model_name, 
        temperature=0, 
        openai_api_key=api_key, 
        response_format={"type": "json_object"},
        max_tokens=4000  # Set a higher token limit
    )
    
    num_slide = num_slides + 5
    # Use a raw string with named parameters
    final_input_prompt = f"""
    You are a professional presentation designer. Create a comprehensive slide deck with {num_slide} pages based on the following user query, context and Guided line:
    
    "{final_prompt}"
    
    Ensure content is clear, concise, and professionally written. Use appropriate slide types for different content needs.
    
    IMPORTANT: 
    1. Return ONLY valid JSON without any additional text, explanations, or formatting.
    2. Ensure exactly {num_slide} slides are included.
    3. Each slide should have substantial content, not just a title.
    4. Maintain consistent detail level across all slides.
    5. Do not reference previous responses or maintain context from previous calls.
    """
    
    # Generate the slide content using the RAG pipeline
    response = llm.invoke(final_input_prompt)

    # Output the generated text for review
    raw_content = response.content

    print(raw_content)
    # Try to parse the entire response as JSON first
    return json.loads(raw_content)



def create_presentation(data, model_name, client, filename="output.pptx", generate_image=False):
    prs = Presentation(ppt_template_path)
    aws_path = st.secrets["AWS_path"]
    image_caption = load_image_captions(aws_path)
    print(data)
    # --- Define Layout Mapping ---
    # !! CRITICAL: Verify these indices match your ACTUAL template 'PPT_01.pptx' !!
    # Layout indices are 0-based. Check Slide Master view in PowerPoint.
    layout_map = {
        "standard": 4,              # Title + Content + Picture Placeholder
        "standard_image": 4,        # Title + Content + Picture Placeholder
        "three_columns": 7,         # Layout specified for three columns
        "progress (Four stage)": 8,              # Layout specified for progress stages
        "comparison": 9,            # Layout specified for comparison (needs verification)
        "process_steps (Three steps)": 10,         # Layout specified for process steps
        "pie_chart": 6              # Reuse a layout with Title + Content (like standard_image)
        # Add other slide_type mappings if needed
    }
    print(f"ℹ️ Layout Mapping Used: {layout_map}")

    # --- Cover Slide ---
    print("\n--- Processing Cover Slide ---")
    try:
        # Use layout 0 (typically Title Slide layout)
        # Note: If template was loaded, prs.slides[0] might be the *first existing slide*
        # It's safer to reference layouts directly if creating from scratch, but here we modify slide 0
        if len(prs.slides) > 0:
             cover_slide = prs.slides[0]
        else:
             print("⚠️ Template contains no slides. Cannot modify cover slide.")

        # Set Title
        title_shape = cover_slide.shapes.title
        title_text = data.get("title", "Presentation Title")
        words = title_text.split()
        new_title = title_text if len(title_text) <= 17 else " ".join(words[:2]) + "\n" + " ".join(words[2:])
        title_shape.text = new_title
        print(f"  - Title set: '{new_title[:50]}...'")


        # Set Subtitle (Assuming Placeholder Index 1 for subtitle on Layout 0)
        subtitle_placeholder_index = 1
        if len(cover_slide.placeholders) > subtitle_placeholder_index:
            cover_slide.placeholders[subtitle_placeholder_index].text = data.get("subtitle", "")
            print(f"  - Subtitle set: '{data.get('subtitle', '')[:50]}...'")
        else:
             print(f"⚠️ Cover slide layout missing subtitle placeholder (expected index {subtitle_placeholder_index}).")


        # Add Cover Image
        cover_image_prompt = data.get("cover_image_prompt")
        if cover_image_prompt and client:
            print(f"  🖼️ Generating cover image...")

            image_caption, cover_image_stream = search_and_generate_images(image_caption, cover_image_prompt, generate_image)

            print(f"  ✔️ Cover image generated.")

        # --- insert photo function ---
        if not generate_image:
            # Simply create a BytesIO object from the image file path
            print(f"  ℹ️ Loading image from path: {cover_image_stream}")
            
            # Create BytesIO object directly from the image file
            # with open(cover_image_stream, 'rb') as img_file:
            #     image_bytes_io = io.BytesIO(img_file.read())
            conn = st.connection('s3', type=FilesConnection)
            # Open the file as bytes
            with conn.open(cover_image_stream, "rb") as f:
                image_bytes = f.read()
            image_bytes_io = io.BytesIO(image_bytes)

            # Rewind the stream to the beginning
            image_bytes_io.seek(0)
            print(f"  ✔️ Image loaded into memory buffer")

            # --- Get the actual pixel dimensions from the image stream ---
            try:
                with Image.open(image_bytes_io) as img:
                    pixel_width, pixel_height = img.size
                    print(f"  ℹ️ Image pixel dimensions: {pixel_width} x {pixel_height}")
                # Rewind stream AGAIN because Image.open consumes it
                image_bytes_io.seek(0)
            except Exception as e:
                print(f"  ❌ Error reading image dimensions with Pillow: {e}")
                # Handle error: maybe set default pixel dimensions or exit
                pixel_width, pixel_height = 100, 100 # Example fallback

            # --- Calculate the size in EMUs based on pixels ---
            # Use a standard DPI, e.g., 96, which is common for screen resolution mapping.
            DPI = 96
            # Convert pixels to inches, then inches to EMUs using pptx.util
            width_in_inches = pixel_width / DPI
            height_in_inches = pixel_height / DPI

            # Use Inches to convert directly to EMU
            width_emu = Inches(width_in_inches)
            height_emu = Inches(height_in_inches)
            # Check if the image width exceeds the maximum allowed width (5245100 EMUs)
            MAX_WIDTH_EMU = 5245100
            

            # Calculate new height while maintaining aspect ratio
            aspect_ratio = height_emu / width_emu
            # Adjust width to maximum allowed
            width_emu = MAX_WIDTH_EMU
            # Recalculate height to maintain aspect ratio
            height_emu = int(width_emu * aspect_ratio)
            print(f"  ⚠️ Image width exceeded maximum allowed. Resizing to width={width_emu}, height={height_emu}")
                
            print(f"  ℹ️ Calculated EMU dimensions: Width={width_emu}, Height={height_emu}")

            # Find Picture placeholder (Type 18) to get its position
            for placeholder in cover_slide.placeholders:
                if placeholder.placeholder_format.type == 18:  # Picture placeholder type
                    # Get the mid-right corner coordinates from the placeholder
                    # X-coordinate of the placeholder's right edge
                    target_left = placeholder.left
                    # Y-coordinate of the placeholder's vertical midpoint
                    target_top = placeholder.top + (placeholder.height / 2) - (height_emu / 2)
                    print(f"  ℹ️ Found placeholder at right: {target_left}, mid: {target_top}")

                    # --- Use add_picture instead of insert_picture ---
                    try:
                        # Add the picture directly to the slide's shapes collection
                        # Position it at the placeholder's top-left
                        # Size it using the calculated EMUs based on original pixel dimensions
                        added_picture = cover_slide.shapes.add_picture(
                            image_bytes_io,
                            target_left,      # Use placeholder's right edge for picture's left
                            target_top,       # Use placeholder's mid-point for picture's top
                            width=width_emu,  # Use calculated image width
                            height=height_emu # Use calculated image height
                        )
                        print(f"  ✔️ Picture added at placeholder position with calculated 'original' size.")
                        # Verify the size set
                        print(f"  ℹ️ Final shape dimensions: Width={added_picture.width}, Height={added_picture.height}")

                        # **Optional:** If you want to remove the original placeholder shape
                        # (which might still be visible underneath or have formatting), you can:
                        placeholder_shape = placeholder.element
                        placeholder_shape.getparent().remove(placeholder_shape)
                        print("  ℹ️ Original placeholder shape removed.")

                    except Exception as e:
                        print(f"  ❌ Error using shapes.add_picture: {e}")

                    break # Exit loop once placeholder is found and picture added
        else:
            for placeholder in cover_slide.placeholders:
                if placeholder.placeholder_format.type == 18:
                    placeholder.insert_picture(cover_image_stream)
                    break


    except Exception as e:
        print(f"❌ Error processing cover slide: {e}")
        traceback.print_exc()

    # --- Content Slides ---
    print("\n--- Processing Content Slides ---")
    slides_data = data.get("slides", [])
    if not isinstance(slides_data, list) or not slides_data:
        print("⚠️ No 'slides' data found or list is empty. Skipping content slides.")
    else:
        # Iterate starting from slide 1 if cover slide was modified, otherwise 0
        # However, we are *adding* new slides, so index doesn't matter here.
        for i, slide_data in enumerate(tqdm(slides_data, desc="📊 Building slides")):
            slide_type = slide_data.get("slide_type", "standard") # Default to standard
            print(f"\nProcessing Slide {i+1}/{len(slides_data)}: Type '{slide_type}'")

            if not isinstance(slide_data, dict):
                print(f"⚠️ Slide data at index {i} is not a dictionary. Skipping.")
                continue

            try:
                # --- Choose Layout ---
                layout_index = layout_map.get(slide_type)

                # if slide_type == "standard" and slide_data.get("image_prompt"):
                #     layout_index = layout_map.get("standard_image")


                if layout_index is None:
                    print(f"  ⚠️ Unknown slide type '{slide_type}' or missing mapping. Using default layout 5.")
                    layout_index = 5 # Default to Title + Content

                try:
                     layout = prs.slide_layouts[layout_index]
                     slide = prs.slides.add_slide(layout)
                     print(f"  - Added slide with Layout Index: {layout_index}")
                except IndexError:
                     print(f"  ❌ Layout index {layout_index} invalid for the template! Check layout_map. Skipping slide.")
                     continue
                except Exception as e:
                     print(f"  ❌ Error adding slide with layout index {layout_index}: {e}. Skipping slide.")
                     continue

                # --- Populate Title ---
                slide.shapes.title.text = slide_data.get("header", f"Slide {i+1}")
                print(f"  - Title set: '{slide_data.get('header', '')[:50]}...'")


                # --- Populate Content based on Slide Type ---
                # Define placeholder indices (these are assumptions based on common layouts)
                # Remember: Placeholder indices (0, 1, 2...) might not match Placeholder IDs (0, 10, 11...)
                main_content_placeholder_index = 1
                placeholder_indices_layout7 = [(13, 10), (14, 11), (15, 12)] # Assumed map for IDs 10, 11, 12 on Layout 7 (three_columns)
                placeholder_indices_layout8_stages = [(14, 11), (15, 18), (16, 19), (17, 20)] # Assumed map for IDs on Layout 8 (progress) very messy...
                placeholder_indices_layout9 = 10 # Assumed map for content on Layout 9 (comparison)
                placeholder_indices_layout10 = [(13, 10), (14, 11), (15, 12)] # Assumed map for BODY IDs 10, 11, 12 on Layout 10 (process_steps)
                placeholder_indices_layout6 = 11 # pie chart ID
                
                # = = = = = = STANDARD = = = = = =
                if slide_type == "standard":
                    if len(slide.placeholders) > main_content_placeholder_index:
                        content_placeholder = slide.placeholders[main_content_placeholder_index]
                        tf = content_placeholder.text_frame
                        tf.clear()
                        bullets = slide_data.get("bullets", [])
                        if not bullets: print("  - No bullets found for standard slide.")
                        for bullet in bullets:
                            if isinstance(bullet, dict):
                                p_main = tf.add_paragraph()
                                p_main.text = bullet.get("text", "")
                                p_main.level = 0

                                p_expl = tf.add_paragraph()
                                p_expl.text = bullet.get("explanation", "")
                                p_expl.level = 1 # Indent explanation
                            else: print(f"  ⚠️ Bullet item is not a dictionary: {bullet}")
                        print(f"  - Populated standard content placeholder {main_content_placeholder_index}.")

                        # Handle optional background image
                        # if is_standard_image and client:
                        image_prompt = slide_data.get("image_prompt")
                        if image_prompt is None:
                            # Generate image prompt using LLM based on slide content
                            llm = ChatOpenAI(model=model_name, temperature=0.7, client=client)
                            slide_content = json.dumps(slide_data)
                            prompt_generation_text = f"Based on this slide content: {slide_content}, generate a concise and descriptive image prompt that would create a relevant background image for this slide. The prompt should be specific and visual."
                            response = llm.invoke(prompt_generation_text)
                            image_prompt = response.content
                        print(f"  🖼️ Generating background image for image prompt: {image_prompt}...")

                        # ... (function call) ...
                        image_caption, image_stream = search_and_generate_images(image_caption, image_prompt, generate_image)
                        print(f"  ✔️ Background image generated.")
                            
                        # --- insert photo function ---
                        if not generate_image:
                            # Simply create a BytesIO object from the image file path
                            print(f"  ℹ️ Loading image from path: {image_stream}")
                            
                            conn = st.connection('s3', type=FilesConnection)
                            with conn.open(image_stream, "rb") as f:
                                image_bytes = f.read()
                            image_bytes_io = io.BytesIO(image_bytes)
                            
                            # Rewind the stream to the beginning
                            image_bytes_io.seek(0)
                            print(f"  ✔️ Image loaded into memory buffer")

                            # --- Get the actual pixel dimensions from the image stream ---
                            try:
                                with Image.open(image_bytes_io) as img:
                                    pixel_width, pixel_height = img.size
                                    print(f"  ℹ️ Image pixel dimensions: {pixel_width} x {pixel_height}")
                                # Rewind stream AGAIN because Image.open consumes it
                                image_bytes_io.seek(0)
                            except Exception as e:
                                print(f"  ❌ Error reading image dimensions with Pillow: {e}")
                                # Handle error: maybe set default pixel dimensions or exit
                                pixel_width, pixel_height = 100, 100 # Example fallback

                            # --- Calculate the size in EMUs based on pixels ---
                            # Use a standard DPI, e.g., 96, which is common for screen resolution mapping.
                            DPI = 96
                            # Convert pixels to inches, then inches to EMUs using pptx.util
                            width_in_inches = pixel_width / DPI
                            height_in_inches = pixel_height / DPI

                            # Use Inches to convert directly to EMU
                            width_emu = Inches(width_in_inches)
                            height_emu = Inches(height_in_inches)
                            # Check if the image width exceeds the maximum allowed width (5245100 EMUs)

                            # Calculate new height while maintaining aspect ratio
                            aspect_ratio = height_emu / width_emu
                            # Adjust width to maximum allowed
                            width_emu = 4892842 # width for placeholder in content
                            # Recalculate height to maintain aspect ratio
                            height_emu = int(width_emu * aspect_ratio)
                            print(f"  ⚠️ Resizing to width={width_emu}, height={height_emu}")
                                
                            print(f"  ℹ️ Calculated EMU dimensions: Width={width_emu}, Height={height_emu}")

                            # Find Picture placeholder (Type 18) to get its position
                            for placeholder in slide.placeholders:
                                if placeholder.placeholder_format.type == 18:  # Picture placeholder type
                                    # Get the mid-right corner coordinates from the placeholder
                                    # X-coordinate of the placeholder's right edge
                                    target_left = placeholder.left
                                    # Y-coordinate of the placeholder's vertical midpoint
                                    target_top = placeholder.top + (placeholder.height / 2) - (height_emu / 2)
                                    print(f"  ℹ️ Found placeholder at right: {target_left}, mid: {target_top}")

                                    # --- Use add_picture instead of insert_picture ---
                                    try:
                                        # Add the picture directly to the slide's shapes collection
                                        # Position it at the placeholder's top-left
                                        # Size it using the calculated EMUs based on original pixel dimensions
                                        added_picture = slide.shapes.add_picture(
                                            image_bytes_io,
                                            target_left,      # Use placeholder's right edge for picture's left
                                            target_top,       # Use placeholder's mid-point for picture's top
                                            width=width_emu,  # Use calculated image width
                                            height=height_emu # Use calculated image height
                                        )
                                        print(f"  ✔️ Picture added at placeholder position with calculated 'original' size.")
                                        # Verify the size set
                                        print(f"  ℹ️ Final shape dimensions: Width={added_picture.width}, Height={added_picture.height}")

                                        # **Optional:** If you want to remove the original placeholder shape
                                        # (which might still be visible underneath or have formatting), you can:
                                        placeholder_shape = placeholder.element
                                        placeholder_shape.getparent().remove(placeholder_shape)
                                        print("  ℹ️ Original placeholder shape removed.")

                                    except Exception as e:
                                        print(f"  ❌ Error using shapes.add_picture: {e}")

                                    break # Exit loop once placeholder is found and picture added
                        else:
                            for placeholder in slide.placeholders:
                                if placeholder.placeholder_format.type == 18:
                                    placeholder.insert_picture(image_stream)
                                    break
                        
                    else:
                        print(f"  ⚠️ Layout {layout_index} for standard slide {i+1} missing content placeholder (expected index {main_content_placeholder_index}).")

                # = = = = = = THREE COLUMNS = = = = = = (finish)
                elif slide_type == "three_columns":
                    columns_data = slide_data.get("columns", [])
                     # Using assumed indices from placeholder_indices_layout7
                    print(f"  - Populating three columns using placeholder indices: {placeholder_indices_layout7}")
                    for idx, col_data in enumerate(columns_data):
                        ph_index = placeholder_indices_layout7[idx]
                        # Get the title and description placeholders from the tuple
                        title_ph_idx, desc_ph_idx = ph_index
                        # Add title to the title placeholder (13, 14, or 15)
                        slide.placeholders[title_ph_idx].text = col_data.get('title', f'Step {idx+1}')
                        # Add description to the description placeholder (10, 11, or 12)
                        column_placeholder = slide.placeholders[desc_ph_idx]
                        c_tf = column_placeholder.text_frame
                        c_tf.clear()

                        points = col_data.get("points", [])
                        for point in points:
                            p_main = c_tf.add_paragraph()
                            p_main.text = point 
                            p_main.level = 0


                # = = = = = = PROGRESS = = = = = = (finish)
                elif slide_type == "progress (Four stage)":
                    stages_data = slide_data.get("stages", [])
                    num_stages = len(stages_data)
                    print(f"  - Found {num_stages} stages for progress slide.")
                    # Using highly speculative mapping from placeholder_indices_layout8_stages

                    for s_idx, stage_data in enumerate(stages_data):
                        # Get the tuple of placeholder indices for this stage
                        stage_ph_indices = placeholder_indices_layout8_stages[s_idx]
                        title_ph_idx, desc_ph_idx = stage_ph_indices
                        print(f"  - Processing Stage {s_idx+1} using placeholder indices: {stage_ph_indices}")

                        # Populate stage title
                        slide.placeholders[title_ph_idx].text = stage_data.get('title', f'Stage {s_idx+1}')

                        # Populate stage description
                        slide.placeholders[desc_ph_idx].text = stage_data.get('description', '')


                # = = = = = = COMPARISON = = = = = =
                elif slide_type == "comparison":
                    # Get the header and items data
                    header = slide_data.get("header", "Comparison")
                    items_data = []
                    
                    # Collect all items for comparison
                    item_index = 1
                    while True:
                        item_key = f"item{item_index}"
                        item_data = slide_data.get(item_key)
                        if item_data is None:
                            break
                        if isinstance(item_data, dict):
                            items_data.append(item_data)
                        item_index += 1
                    
                    # If no items found using sequential naming, try to get items as a list
                    if not items_data:
                        items_list = slide_data.get("items", [])
                        if isinstance(items_list, list):
                            items_data = [item for item in items_list if isinstance(item, dict)]
                    
                    num_items = len(items_data)
                    
                    # Using assumed indices from placeholder_indices_layout9
                    if num_items >= 2:  # Need at least 2 items for comparison
                        print(f"  - Creating comparison table for slide {i+1} with header: '{header}' and {num_items} items")

                        # Get the main content placeholder
                        try:
                            content_placeholder = slide.placeholders[placeholder_indices_layout9]
                        except (KeyError, IndexError):
                            print(f"  ⚠️ Content placeholder index {placeholder_indices_layout9} not found in slide {i+1}")
                            # Try to find a suitable content placeholder
                            content_placeholder = None
                            for ph_idx, ph in enumerate(slide.placeholders):
                                if ph.placeholder_format.type in [2, 7, 12]:  # Content or body placeholder types
                                    content_placeholder = ph
                                    print(f"  - Found alternative content placeholder at index {ph_idx}")
                                    break
                            
                            if content_placeholder is None:
                                print(f"  ❌ No suitable content placeholder found for slide {i+1}. Skipping table creation.")
                                continue
                        
                        # Get points from all items
                        all_item_points = [item_data.get("points", []) for item_data in items_data]
                        
                        # Create a list of metrics based on the points from all items
                        metrics = slide_data.get('metrics', [])
                        if not metrics:
                            max_points = max(len(points) for points in all_item_points)
                            for idx in range(max_points):
                                 metrics.append(f"Metric {idx+1}")
                        
                        # Create table with rows for header + each point and columns for metric + each item
                        rows = len(metrics) + 1  # +1 for the header row
                        cols = num_items + 1  # +1 for the metric column
                        
                        # Add table to the slide
                        table = content_placeholder.insert_table(rows, cols).table
                        
                        # Set column widths proportionally
                        total_width = sum(col.width for col in table.columns)
                        metric_width_pct = 0.15  # 15% for metric column
                        item_width_pct = (1 - metric_width_pct) / num_items  # Distribute remaining width evenly
                        
                        table.columns[0].width = int(total_width * metric_width_pct)
                        for col_idx in range(1, cols):
                            table.columns[col_idx].width = int(total_width * item_width_pct)
                        
                        # Set header row with item titles
                        for col_idx, item_data in enumerate(items_data):
                            table.cell(0, col_idx + 1).text = item_data.get("title", f"Item {col_idx + 1}")
                        
                        # Fill in the comparison points row by row
                        for row_idx, metric in enumerate(metrics):
                            # Set the metric name in first column
                            table.cell(row_idx + 1, 0).text = metric
                            
                            # Set the value for each item if available
                            for col_idx, item_points in enumerate(all_item_points):
                                if row_idx < len(item_points):
                                    point = item_points[row_idx]
                                    # Extract value (after the colon)
                                    value = point.split(':', 1)[1].strip() if ':' in point else point
                                    table.cell(row_idx + 1, col_idx + 1).text = value
                        
                        # Format the table for better appearance
                        for row_idx in range(rows):
                            for col_idx in range(cols):
                                cell = table.cell(row_idx, col_idx)
                                paragraph = cell.text_frame.paragraphs[0]
                                
                                # Center align all cells
                                paragraph.alignment = PP_ALIGN.CENTER
                                
                                # Apply font formatting
                                for run in paragraph.runs:
                                    run.font.size = Pt(20)  # Adjust font size as needed
                                    
                                    # Make header row bold with white text, and header column bold with dark gray text
                                    if row_idx == 0:
                                        run.font.bold = True
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                                    elif col_idx == 0:
                                        run.font.bold = True
                                        run.font.color.rgb = RGBColor(64, 64, 64)  # Dark gray for header column
                                
                                # Apply cell formatting for header row
                                if row_idx == 0:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(148, 189, 94)  # light greem header
                                
                                # Apply alternating row colors for better readability
                                elif row_idx % 2 == 1:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Light gray
                        
                        print(f"  ✓ Successfully created comparison table with {rows} rows and {cols} columns")
                    
                    else:
                        print(f"  ⚠️ Missing or invalid data for comparison slide {i+1}. Expected 'item1' and 'item2' dictionaries.")
                # = = = = = = PROCESS STEPS = = = = = =(finish)
                elif slide_type == "process_steps (Three steps)":
                    steps_data = slide_data.get("steps", [])
                    num_steps = len(steps_data)
                    # Using assumed indices from placeholder_indices_layout10
                    max_steps_layout = len(placeholder_indices_layout10) # e.g., 3 steps

                    if len(slide.placeholders) >= max_steps_layout + 1 : # Need at least title + step placeholders
                        print(f"  - Populating process steps using placeholder indices: {placeholder_indices_layout10}")
                        if num_steps == 0: print("  - No steps found in data.")
                        for s_idx, step_data in enumerate(steps_data):
                            if s_idx >= max_steps_layout:
                                print(f"  ⚠️ More steps ({num_steps}) than layout mapping supports ({max_steps_layout}). Skipping extra.")
                                break

                            if not isinstance(step_data, dict):
                                 print(f"  ⚠️ Step data {s_idx} is not a dictionary.")
                                 continue

                            ph_index = placeholder_indices_layout10[s_idx]
                            try:
                                 # Get the title and description placeholders from the tuple
                                 title_ph_idx, desc_ph_idx = ph_index
                                 # Add title to the title placeholder (13, 14, or 15)
                                 slide.placeholders[title_ph_idx].text = f"{s_idx+1:02d}\n{step_data.get('title', f'Step {s_idx+1}')}"
                                 # Add description to the description placeholder (10, 11, or 12)
                                 slide.placeholders[desc_ph_idx].text = step_data.get("description", "")

                            except IndexError:
                                print(f"  ❌ Error accessing placeholder index {ph_index} for step {s_idx+1}. Check Layout 10 structure.")
                                break # Stop if layout is wrong
                            except Exception as e:
                                print(f"  ❌ Error populating step {s_idx+1} using index {ph_index}: {e}")
                    else:
                          print(f"  ⚠️ Layout {layout_index} for process_steps slide {i+1} does not have enough text placeholders (expected indices {placeholder_indices_layout10}). Check Layout 10 structure.")

                # = = = = = = PIE CHART = = = = = =(finish)
                elif slide_type == "pie_chart":
                    chart_data = CategoryChartData()
                    data_points = slide_data.get("data_points", [])
                    chart_description = slide_data.get("data_description", "Data Distribution") # Use description as chart title/subtitle
                    print(f"  - Processing Pie Chart: {chart_description}")

                    # Prepare chart data
                    categories = []
                    values = []
                    unit = slide_data.get("data_unit", "")  # Get unit from slide data if available
                    for i, dp in enumerate(data_points):
                        label = dp.get("label", f"Category {i+1}")
                        value = dp.get("value")
                        categories.append(label)
                        values.append(float(value))


                    chart_data.categories = categories
                    chart_data.add_series("", values) # Series name usually not shown on pie
                    print(f"  - Chart data prepared with {len(categories)} categories.")

                    # Add chart to slide
                    # Prefer inserting into placeholder (assume index 1 for layout 4)
                    chart_placeholder = slide.placeholders[placeholder_indices_layout6]
                    graphic_frame = chart_placeholder.insert_chart(
                        XL_CHART_TYPE.PIE, chart_data
                    )
                    chart = graphic_frame.chart

                    print(f"  📊 Pie chart inserted into placeholder index {placeholder_indices_layout6}.")
                    
                    # Customize chart
                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.RIGHT  # Use enum value instead of integer
                    chart.legend.include_in_layout = False
                    chart.legend.font.size = Pt(20)

                    plot = chart.plots[0]
                    plot.has_data_labels = True
                    data_labels = plot.data_labels
                    
                    # Apply number format with unit
                    if unit.lower() == "percent" or unit.lower() == "%":
                        data_labels.number_format = '0 "%"'  # Show as 70%
                    elif unit:
                        data_labels.number_format = f'0 "{unit}"'  # Show as 70 unit
                    else:
                        data_labels.number_format = '0'  # Show as 70
                        
                    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                    data_labels.font.size = Pt(15)
                    
                    # Add chart description
                    content_placeholder = slide.placeholders[main_content_placeholder_index]
                    tf = content_placeholder.text_frame
                    tf.clear()
                    p_main = tf.add_paragraph()
                    p_main.text = chart_description
                    p_main.level = 0


                # = = = = = = OTHER / UNKNOWN = = = = = =
                else:
                    print(f"  ℹ️ Slide type '{slide_type}' not explicitly handled. Check JSON or add logic.")
                    # Optional: Add default text if a content placeholder exists
                    if len(slide.placeholders) > main_content_placeholder_index:
                        try:
                            slide.placeholders[main_content_placeholder_index].text = f"Content for slide type: {slide_type}\nData: {str(slide_data)[:100]}..."
                        except Exception: pass

            except Exception as e:
                print(f"❌❌❌ Major error processing slide {i+1} (Type: {slide_type}): {e}")
                traceback.print_exc() # Print detailed traceback
                # continue # Continue to next slide even if one fails

    # --- Save Presentation ---
    print("\n--- Saving Presentation ---")
    try:
        prs.save(filename)
        print(f"✅ Presentation saved successfully to: {filename}")
    except Exception as e:
        print(f"❌ Error saving presentation to '{filename}': {e}")
        print("  - Check file permissions and if the file is open elsewhere.")


def revise_slides(slides_data, revision_instruction, api_key = None, model_name="gpt-4o-mini-2024-07-18"):
    """
    Uses the LLM model to revise slide content based on user instructions,
    reusing the same memory context from the initial generation.
    
    Args:
        slides_data (dict): The JSON structure containing the presentation data
        user_query (str): The original presentation topic query
        rag (RAGModel, optional): The RAG model from the initial generation
        memory (ConversationBufferMemory, optional): Memory from initial generation
        num_slides (int): Target number of slides
        
    Returns:
        dict: The revised slides data
    """

    # Convert slides_data to a formatted string for the prompt
    slides_json = json.dumps(slides_data, indent=2)
    slides_json = slides_json.replace("{", "{{")
    slides_json = slides_json.replace("}", "}}")
    
    # Create the revision prompt using RAG
    revision_prompt = f"""
You are a professional presentation designer tasked with revising an existing presentation based on the user's request.

Here is the ORIGINAL presentation structure in JSON format:
{slides_json}

IMPORTANT - REVISION REQUEST: "{revision_instruction}"

Focus specifically on implementing the changes requested in the revision instruction above. Make targeted modifications to the presentation structure to address these specific requests while maintaining the overall quality and coherence of the presentation.

If the revision requests a change in the number of slides, please note that the cover page counts as one slide. For example, if the request is for "10 slides total", you should include the cover page plus 9 content slides. Adjust the number of slides in the JSON structure accordingly while preserving the most important content.
 
 Return the COMPLETE revised JSON structure with all fields, ensuring that all requested changes have been properly implemented.
 
 Include in your JSON:
 - "title": A concise, impactful presentation title (4-6 words)
 - "subtitle": A clarifying subtitle that provides context (10-15 words)
 - "cover_image_prompt": Detailed DALL-E prompt for a professional, relevant cover image (approx. 100-150 chars). Must be safe for work."
 - "slides": An array of slide objects, each potentially having a different structure based on `slide_type`:
 
 **JSON Structure Required:**

{{{{
  "title": "Concise, impactful presentation title (5-7 words)",
  "subtitle": "Clarifying subtitle providing context (10-15 words)",
  "cover_image_prompt": "Detailed DALL-E prompt for a professional, relevant cover image (approx. 100-150 chars). Must be safe for work.",
  "slides": [
    // --- Example Slide Type: Standard (Bullet Points) ---
    {{{{
      "slide_type": "standard", // Type: standard, pie_chart, progress, comparison, process_steps, three_columns etc.
      "header": "Engaging Slide Title (max 10 words)",
      "bullets": [ // Use 'bullets' only for 'standard' type slides
        {{{{
          "text": "Clear, actionable key point (5-10 words)",
          "explanation": "Insightful explanation adding value (2-3 complete sentences)"
        }}}}
        // ... 3-4 bullet points total ...
      ],
      "include_background_image": true, // boolean: true if a background/thematic image enhances this slide
      "image_prompt": "Optional: Detailed DALL-E prompt for a relevant, professional BACKGROUND image (if include_background_image is true, otherwise empty string). Must be safe for work."
    }}}},
    // --- Example Slide Type: Pie Chart ---
    {{{{
      "slide_type": "pie_chart",
      "header": "Data Distribution / Market Share", // Title for the chart slide
      "data_description": "Complete summary of what the pie chart illustrates (1-2 sentences explaining the data source, significance, and key insights). For example: 'Q1 Sales Distribution by Region showing the Northeast leading with 45% market share' or 'Resource Allocation Percentages highlighting our focus on R&D investment'.", // Context for the data that will appear as a caption or explanation
      "data_unit": "The unit of measurement for the data (e.g., '%', '$', 'units', 'people'). This will be used for formatting labels.",
      "data_points": [ // An array representing the slices of the pie
        {{{{
          "label": "Category A", // Label for the first slice
          "value": 45 // Numerical value for the first slice (e.g., percentage or amount)
        }}}},
        {{{{
          "label": "Category B",
          "value": 30
        }}}},
        {{{{
          "label": "Category C",
          "value": 25
        }}}}
        // ... Add more data points as provided in the context ...
      ],
    }}}},
    // --- Example Slide Type: Progress (Four Stage-Based) ---
    {{{{
      "slide_type": "progress (Four stage)",
      "header": "Project Milestones & Current Status", // Title for the progress slide
      "stages": [ // An array representing the sequential stages/milestones
          // ... Four stage objects with title and description
      ],
    }}}},
    // --- Example Slide Type: Comparison (Two Items) ---
    {{{{
      "slide_type": "comparison",
      "header": "Comparison: [Item 1 Name] vs. [Item 2 Name]", // Main title for the comparison
      "metrics": [
        "Cost Efficiency", 
        "Performance", 
        "Scalability", 
        "User Experience", 
        "Implementation Time"
      ], // At least 5 metrics/categories for comparison
      "item1": {{{{
        "title": "Item 1 Title",
        "points": [
          "Cost analysis: initial investment, operational expenses, and long-term financial implications.",
          "Performance evaluation: speed, reliability, and throughput capacity.",
          "Scalability: handling increasing workloads, resource requirements, and expansion capabilities.",
          "User experience: interface design, accessibility features, and learning curve.",
          "Implementation: timeline, resource requirements, training needs, and integration challenges."
        ]
      }}}},
      "item2": {{{{
        "title": "Item 2 Title",
        "points": [
          "Cost analysis: initial investment, operational expenses, and long-term financial implications.",
          "Performance evaluation: speed, reliability, and throughput capacity.",
          "Scalability: handling increasing workloads, resource requirements, and expansion capabilities.",
          "User experience: interface design, accessibility features, and learning curve.",
          "Implementation: timeline, resource requirements, training needs, and integration challenges."
        ]
      }}}}
    }}}},
    // --- Example Slide Type: Process Steps (Three steps) ---
    {{{{
      "slide_type": "process_steps (Three steps)",
      "header": "Our [Number]-Step Workflow", // Title for the process slide
      "steps": [ // An array representing the sequential steps
          // ... step objects with title and description ...
      ],
    }}}},
    // --- Example Slide Type: Three Columns ---
    {{{{
      "slide_type": "three_columns",
      "header": "Key Focus Areas", // Title for the slide
      "columns": [ // An array containing exactly three objects, one for each column
         // ... column objects with title and points ...
      ],
    }}}}
    // ... adapting slide_type and content fields as needed ...
  ]
}}}}

**Presentation Best Practices to Follow:**

- **Diversify slide types:** Choose from a variety of slide formats for engagement.
- **Balance visual and text:** Use visual slides (pie_chart, progress (Four stage), three_columns, process_steps (Three steps), comparison) to break up text-heavy sections. Aim for 60% visual, 40% text ratio.
- **Data presentation:** When using charts, ensure data is meaningful and properly contextualized with a clear `data_description`.
- **Narrative structure:** Begin with a compelling introduction, develop key points in logical sequence, and end with impactful conclusion or action items.
- **Content hierarchy:** Prioritize information with main points first, supporting details second.
- **Engagement elements:** Include interactive components, questions, or discussion prompts where appropriate.
- **Visual consistency:** Maintain cohesive design elements throughout. Image prompts should align with slide content and overall presentation theme.

**Output Requirements:**
- Replace any occurrence of "Capital Area Food Bank" in the title with "CAFB".
- Return *only* the valid JSON object.
- Do not include any introductory text, explanations, apologies, or markdown formatting like ```json before or after the JSON.
- Ensure all strings within the JSON are properly escaped.
"""
    # Use a raw string with named parameters
    llm = ChatOpenAI(
        model=model_name, 
        temperature=0, 
        openai_api_key=api_key, 
        response_format={"type": "json_object"},
        max_tokens=4000  # Set a higher token limit
    )
    
    try:
        
        # Use the RAG model to generate the revised slides
        response = llm.invoke(revision_prompt)
        # Extract JSON from the response
        raw_content = response.content
        match = re.search(r'\{.*\}', raw_content, re.DOTALL)
        
        if match:
            revised_slides = json.loads(match.group())
            print("✅ Slides successfully revised!")
            return revised_slides
        else:
            print("❌ Could not extract valid JSON from the response. Keeping original slides.")
            return slides_data
        
    except Exception as e:
        print(f"❌ Error during revision: {e}")
        print("Continuing with original slides.")
        return slides_data

def convert_pptx_to_images_mac(pptx_path):
    """Convert PowerPoint slides to a list of images using LibreOffice and pdf2image"""
    from pdf2image import convert_from_path
    # First convert PPTX to PDF using LibreOffice
    # Sanitize filename to avoid special characters that cause path issues
    base_dir = os.path.dirname(pptx_path)
    base_name = os.path.basename(pptx_path).replace(':', '_').replace('/', '_').replace('\\', '_')
    sanitized_pptx_path = os.path.join(base_dir, base_name)
    
    # If the path changed due to sanitization, create a copy with the sanitized name
    if sanitized_pptx_path != pptx_path:
        import shutil
        shutil.copy2(pptx_path, sanitized_pptx_path)
        pptx_path = sanitized_pptx_path
    
    pdf_path = pptx_path.replace('.pptx', '.pdf')

    try:
        # Convert PPTX to PDF using LibreOffice
        cmd = [
            'soffice',
            '--headless',
            '--convert-to',
            'pdf',
            '--outdir',
            os.path.dirname(pptx_path),
            pptx_path
        ]
        subprocess.run(cmd, check=True)

        # Verify PDF was created
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF conversion failed. Output file not found: {pdf_path}")
        
        # Convert PDF to images
        images = convert_from_path(pdf_path)

        # Save images to temporary files
        image_paths = []
        for i, image in enumerate(images):
            temp_path = os.path.join(tempfile.gettempdir(), f'slide_{i}.png')
            image.save(temp_path, 'PNG')
            image_paths.append(temp_path)

        # Clean up files
        os.remove(pdf_path)
        if sanitized_pptx_path != pptx_path and os.path.exists(sanitized_pptx_path):
            os.remove(sanitized_pptx_path)

        return image_paths

    except Exception as e:
        # Clean up files
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        if sanitized_pptx_path != pptx_path and os.path.exists(sanitized_pptx_path):
            os.remove(sanitized_pptx_path)
        raise Exception(f"Error converting presentation: {str(e)}")



def convert_pptx_to_images_win(pptx_path):
    import win32com.client
    import pythoncom
    """Convert PowerPoint slides to a list of images using win32com"""
    # Get absolute path
    abs_path = os.path.abspath(pptx_path)
    output_dir = os.path.dirname(abs_path)
    
    # Initialize COM objects
    pythoncom.CoInitialize()
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.Visible = True
    
    try:
        # Open the presentation
        presentation = powerpoint.Presentations.Open(abs_path)
        
        # Create a list to store image paths
        image_paths = []
        
        # Save each slide as an image
        for i in range(1, presentation.Slides.Count + 1):
            image_path = os.path.join(output_dir, f'slide_{i}.png')
            presentation.Slides(i).Export(image_path, "PNG")
            image_paths.append(image_path)
            
        # Close presentation and quit PowerPoint
        presentation.Close()
        powerpoint.Quit()
        
        return image_paths
        
    except Exception as e:
        if 'presentation' in locals():
            presentation.Close()
        if 'powerpoint' in locals():
            powerpoint.Quit()
        raise e
    finally:
        pythoncom.CoUninitialize()
